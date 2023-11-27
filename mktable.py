import json
from pptx import Presentation
from pptx.oxml.xmlchemy import OxmlElement
import os
import sys

def _SubElement(parent, tagname, **kwargs):
    """Helper for Paragraph bullet Point
    """
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element

def _replace_paragraph_text(p, t, font=None):
    if p.runs:
        p.runs[0].text = t # Replace 'New Text' with your text
        for other_run in p.runs[1:]:
            other_run.text = ''
    else:
        newrun = p.add_run()
        newrun.text = t # Add new run if no runs are present
        if font: # if has font, set it
            newrun.font.name = font.name 
            newrun.font.size = font.size 

def write_cell(cell, text_list):
    if cell.text_frame.paragraphs:
        # If there is already text in the cell
        first_paragraph = cell.text_frame.paragraphs[0]
        for other_paragraph in cell.text_frame.paragraphs[1:]:
            other_paragraph.clear()
    else:
        first_paragraph = cell.text_frame.add_paragraph()


    first_pPr = first_paragraph._p.get_or_add_pPr()
    # can't access buSzPct & buChar, so overwrite them
    _SubElement(parent=first_pPr, tagname="a:buSzPct", val="111000")
    _SubElement(parent=first_pPr, tagname='a:buChar', char="•")
    # Write the first run's text (replace)
    _replace_paragraph_text(first_paragraph, text_list[0])
    org_font = first_paragraph.runs[0].font if first_paragraph.runs else None

    for t in text_list[1:]:
        p = cell.text_frame.add_paragraph()
        p.level = first_paragraph.level
        pPr = p._p.get_or_add_pPr()
        # Set bullet the same as first paragraph
        pPr.set('marL', first_pPr.attrib['marL'])
        pPr.set('indent', first_pPr.attrib['indent'])
        _SubElement(parent=pPr, tagname="a:buSzPct", val="111000")
        _SubElement(parent=pPr, tagname='a:buChar', char="•")
        # Write the new text
        _replace_paragraph_text(p, t, org_font)


def create_table(page):
    ###########################################
    # Create table
    ###########################################

    # Get the path of the current file (__file__)
    current_file_path = os.path.abspath(__file__)
    # Get the directory containing the current file
    script_folder = os.path.dirname(current_file_path)

    template_path = '{}/template/table_page.pptx'.format(script_folder)
    # Load the template
    ppt = Presentation(template_path)

    # Identify columns number
    number_columns = len(page['content'])
    # Check if the column number is valid
    if number_columns < 2 or page_number > 6:
        print("Invalid column number: {number_columns}")
        sys.exit(1)

    # Page 1: 2 columns, Page 2: 3 columns, etc. Page 5: 6 columns
    slide = ppt.slides[number_columns - 2]

    # For debug to print out placeholder's index and name
    # for shape in slide.placeholders:
    #     print(shape.placeholder_format.idx, ":", shape.name)

    # Set the title for the new slide
    slide.shapes.title.text = page['pageTitle']

    for shape in slide.shapes:
        # Check if shape is a table
        if shape.shape_type == 19:  # 19 corresponds to a table
            table = shape.table
            # Iterate over each main point
            for (idx, item) in enumerate(page['content']):
                p = table.cell(0, idx)
                write_cell(p, [item['title']])

                # Iterate over each subpoint
                cell = table.cell(1, idx)
                write_cell(cell, item['subpoints'])

    # Save the modified presentation
    ppt.save('sample_modified_table.pptx')

if __name__ == "__main__":
    # Check if a page number is provided
    if len(sys.argv) < 2:
        print("Please provide a page number.")
        sys.exit(1)
    page_number = int(sys.argv[1]) - 1  # Adjust for zero-based indexing

    # Read and parse the JSON file
    json_filename = 'content.json'
    with open(json_filename, 'r') as file:
        data = json.load(file)
    # Check if the page number is valid
    if page_number >= len(data["pages"]) or page_number < 0:
        print("Invalid page number.")
        sys.exit(1)

    page = data["pages"][page_number]
        
    create_table(page)
    