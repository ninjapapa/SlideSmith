from pptx import Presentation
import sys


# read in file name
ppt_name = sys.argv[1]

# Load the presentation
ppt = Presentation(ppt_name)

# Access a specific slide by index (e.g., the first slide)
slide = ppt.slides[0]

# Iterate through all shapes on the slide
for shape in slide.shapes:
    # Check if shape is a text frame (regular shape or placeholder)
    print("Name:", shape.name)
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print(run.text)

    # Check if shape is a table
    elif shape.shape_type == 19:  # 19 corresponds to a table
        table = shape.table
        num_columns = len(table.columns)
        print(f"Number of Columns: {num_columns}")

        print(dir(table._graphic_frame))
        # Create a new table with an additional column
        # num_rows = len(table.rows)
        # num_cols = 5
        # (left, top, width, height) = (table._graphic_frame.left, table._graphic_frame.top, table._graphic_frame.width, table._graphic_frame.height)
        # new_table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

        # # Copy contents from old table to new table
        # for row_idx in range(num_rows):
        #     for col_idx in range(num_cols):
        #         new_table.cell(row_idx, col_idx).text = "aaaa"

        # slide.shapes._spTree.remove(shape._element)
        for (idx, row) in enumerate(table.rows):
            print("row:", dir(row))
            for (c_idx, cell) in enumerate(row.cells):
                print("column:", c_idx)
                if cell.text:
                    print(cell.text)
                first_paragraph = cell.text_frame.paragraphs[0]
                first_paragraph.runs[0].text = 'hahahaha\n dadsa'  # Replace 'New Text' with your text

    # Handling other shapes like images or charts can be done here
    # For example, shape.shape_type == 13 corresponds to a chart
    # shape.shape_type == 13 corresponds to a picture (image)
    ppt.save('sample_modified.pptx')
