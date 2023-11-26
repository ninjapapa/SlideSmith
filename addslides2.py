import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Get the path of the current file (__file__)
current_file_path = os.path.abspath(__file__)
# Get the directory containing the current file
script_folder = os.path.dirname(current_file_path)

template_path = '{}/template/template1.pptx'.format(script_folder)
# Load the template
ppt = Presentation(template_path)

def add_slides_from_json(json_filename):

    # Read and parse the JSON file
    with open(json_filename, 'r') as file:
        data = json.load(file)

    # Iterate over each page in the JSON
    for page in data['pages']:
        # Define the slide layout
        slide_layout = ppt.slide_layouts[2]  # table

        # Add a slide
        slide = ppt.slides.add_slide(slide_layout)

        # For debug to print out placeholder's index and name
        # for shape in slide.placeholders:
        #     print(shape.placeholder_format.idx, ":", shape.name)

        # Set the title for the new slide
        slide.shapes.title.text = page['pageTitle']

        # Add content to the slide
        content = slide.placeholders[21]
        cols_cnt = len(page['content'])
        shape = content.insert_table(rows=2, cols=cols_cnt)
        table = shape.table
        
        # Iterate over each main point
        for (idx, item) in enumerate(page['content']['list']):
            p = table.cell(0, idx)
            p.text = item['title']

            # Iterate over each subpoint
            for subpoint in item['subpoints']:
                tf = table.cell(1, idx).text_frame
                sp = tf.add_paragraph()
                sp.text = subpoint
                sp.level = 0

    # Save the modified presentation
    ppt.save('sample_modified.pptx')

# Example usage
add_slides_from_json('content.json')

