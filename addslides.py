import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
import os

# Get the path of the current file (__file__)
current_file_path = os.path.abspath(__file__)
# Get the directory containing the current file
script_folder = os.path.dirname(current_file_path)

template_path = '{}/template/template1.pptx'.format(script_folder)
# Load the template
ppt = Presentation(template_path)

# Get width and height of the slide
slide_width = ppt.slide_width
slide_height = ppt.slide_height

def add_slides_from_json(json_filename):

    # Read and parse the JSON file
    with open(json_filename, 'r') as file:
        data = json.load(file)

    # Iterate over each page in the JSON
    for page in data['pages']:
        # Define the slide layout
        slide_layout = ppt.slide_layouts[1]  # Title and content layout

        # Add a slide
        slide = ppt.slides.add_slide(slide_layout)

        # For debug to print out placeholder's index and name
        # for shape in slide.placeholders:
        #     print(shape.placeholder_format.idx, ":", shape.name)

        # Set the title for the new slide
        slide.shapes.title.text = page['pageTitle']

        # Add content to the slide
        content = slide.placeholders[21]  # Assuming the placeholder[10] is for content, user created placeholders are from 10
        tf = content.text_frame

        # Iterate over each main point
        paragraph_cnt = 0
        for item in page['content']:
            if (paragraph_cnt == 0):
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            paragraph_cnt = paragraph_cnt + 1

            p.text = item['title']
            p.level = 0

            # Iterate over each subpoint
            for subpoint in item['subpoints']:
                sp = tf.add_paragraph()
                sp.text = subpoint
                sp.level = 1

        # Define the size and position of the text box
        left = slide_width - Inches(2)  # Halfway across the width of a standard slide
        top = Inches(0)
        width = Inches(2)   # About 1/4 of the page width
        height = Inches(3)  # About 1/2 of the page height

        # Add the text box
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        # Set word wrap
        text_frame.word_wrap = True
        # Set the background color of the text box to golden
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = RGBColor(255, 215, 0)  # RGB values for golden color

        # Add the notes text
        p = text_frame.paragraphs[0]
        p.text = page['notes']
        p.font.size = Pt(12)  # Adjust font size as needed
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Save the modified presentation
    ppt.save('sample_modified.pptx')

# Example usage
add_slides_from_json('content.json')

